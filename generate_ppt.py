from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Vedāntic Study Archive"
subtitle.text = "A repository of Vedāntic study materials\nFeaturing Bṛhadāraṇyaka Upaniṣad with Śāṅkara Bhāṣya"

# Slide 2: Platform Overview
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Platform Overview"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame

p = tf.add_paragraph()
p.text = "Purpose & Mission"
p.level = 0
p = tf.add_paragraph()
p.text = "A dedicated platform hosting detailed class notes from the teachings of Pūjya Swami Shankarananda."
p.level = 1
p = tf.add_paragraph()
p.text = "Focuses on providing deep, textual study aids for the principal Upaniṣads."
p.level = 1

p = tf.add_paragraph()
p.text = "Target Audience"
p.level = 0
p = tf.add_paragraph()
p.text = "Fellow seekers, students of Advaita Vedānta, and scholars interested in the traditional commentaries (Śāṅkara Bhāṣya and Vārttika)."
p.level = 1

# Slide 3: Key Features
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Key Features"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame

p = tf.add_paragraph()
p.text = "1. Bilingual Script Support"
p.level = 0
p = tf.add_paragraph()
p.text = "Seamlessly toggle between Devanagari script and IAST (Roman with diacritics) for every class."
p.level = 1

p = tf.add_paragraph()
p.text = "2. High-Fidelity Transcripts"
p.level = 0
p = tf.add_paragraph()
p.text = "Notes are meticulously prepared from video recordings and reviewed against original classes."
p.level = 1

p = tf.add_paragraph()
p.text = "3. Powerful Search Engine"
p.level = 0
p = tf.add_paragraph()
p.text = "Lightning-fast, full-text search across all classes, scripts, and summaries."
p.level = 1

# Slide 4: Current Series
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Current Series: Bṛhadāraṇyaka Upaniṣad"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame

p = tf.add_paragraph()
p.text = "Structure & Organization"
p.level = 0
p = tf.add_paragraph()
p.text = "Organized hierarchically by Adhyāya and Brāhmaṇa (e.g., Maitreyī Brāhmaṇa)."
p.level = 1

p = tf.add_paragraph()
p.text = "Rich Metadata"
p.level = 0
p = tf.add_paragraph()
p.text = "Each class entry includes concise summaries of the topics covered."
p.level = 1
p = tf.add_paragraph()
p.text = "Direct links to source texts (Advaita Sharada), video playlists, and audio materials."
p.level = 1

# Slide 5: Technical Excellence
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Technical Excellence"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame

p = tf.add_paragraph()
p.text = "Lightning Fast & Accessible"
p.level = 0
p = tf.add_paragraph()
p.text = "Built on Astro (Zero-JS by default) ensuring maximum speed and SEO optimization."
p.level = 1

p = tf.add_paragraph()
p.text = "Responsive & Beautiful"
p.level = 0
p = tf.add_paragraph()
p.text = "Fully responsive design with built-in Light and Dark themes."
p.level = 1
p = tf.add_paragraph()
p.text = "Elegant typography specifically chosen for reading Sanskrit and English."
p.level = 1

# Save the presentation
prs.save("public/Vedantic_Archive_Overview.pptx")
print("Presentation saved to public/Vedantic_Archive_Overview.pptx")
